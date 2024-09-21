import eyed3
import imageio
from fastapi import FastAPI, File, UploadFile
from pptx import Presentation
import requests
import os
from moviepy.editor import VideoFileClip, AudioFileClip, concatenate_videoclips
from pydantic import BaseModel
import shutil
from firstordermodel import demo
from skimage.transform import resize
import matplotlib.pyplot as plt
import matplotlib.animation as animation
from pydub import AudioSegment
from moviepy.video.io.ffmpeg_tools import ffmpeg_merge_video_audio
import cv2
import numpy as np
from moviepy.editor import VideoFileClip, VideoClip, CompositeVideoClip

app = FastAPI()

XI_API_KEY = "38b0fb5c6fe325393ffc6f0883f73c2c"
output_dir = "./outputaudio"
voice_id = "29vD33N1CtxCmqQRPOHJ"

generator, kp_detector = demo.load_checkpoints(config_path='./vox-256.yaml',
                            checkpoint_path='./vox-adv-cpk.pth.tar')
reader = imageio.get_reader('./Obama.mp4')


class DeepFakeRequest(BaseModel):
    ppt_path: str
    image_path: str
    video_path: str
    educator_id: str


# Helper function to extract speaker notes
def extract_speaker_notes(ppt_path):
    presentation = Presentation(ppt_path)
    notes = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            note_text = notes_slide.notes_text_frame.text
            notes.append({
                'slide_number': slide_number,
                'speaker_notes': note_text
            })
    return notes


# Helper function to convert speaker notes to audio using Eleven Labs API
def comment_to_audio(text, i):
    tts_url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}/stream"
    headers = {
        "Accept": "application/json",
        "xi-api-key": XI_API_KEY
    }
    data = {
        "text": text,
        "model_id": "eleven_multilingual_v2",
        "voice_settings": {
            "stability": 0.5,
            "similarity_boost": 0.8
        }
    }
    response = requests.post(tts_url, headers=headers, json=data, stream=True)
    if response.ok:
        output_path = os.path.join(output_dir, f"output{i}.mp3")
        with open(output_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=1024):
                f.write(chunk)
        return output_path
    else:
        return None


# API to upload a PPT and extract speaker notes
def save_generated_video(generated, filename='./genratedvideo/predicted_video.mp4'):
    fig = plt.figure(figsize=(8, 6))

    ims = []
    for frame in generated:
        im = plt.imshow(frame, animated=True)
        plt.axis('off')
        ims.append([im])

    ani = animation.ArtistAnimation(fig, ims, interval=50, repeat_delay=1000)
    ani.save(filename, writer='ffmpeg', fps=30)  # Adjust fps as needed
    plt.close()

def audio_length(filepath):
  return eyed3.load(filepath).info.time_secs

def combine_audios(audio_files, output_file):
    combined_audio = None
    for audio_file in audio_files:
        audio = AudioSegment.from_file(audio_file)
        if combined_audio is None:
            combined_audio = audio
        else:
            combined_audio += audio  # Accumulate each audio file
    print(combined_audio)

    combined_audio.export(output_file, format="mp3")

def repeat_video(input_video, repetitions, output_video):
    clip = VideoFileClip(input_video)

    # Get the FPS from the original clip
    fps = clip.fps

    # Create a list of repeated clips
    repeated_clips = [clip] * int(repetitions)

    # Handle fractional repetitions by appending part of the clip
    remainder_clip_duration = (repetitions - int(repetitions)) * clip.duration
    if remainder_clip_duration > 0:
        repeated_clips.append(clip.subclip(0, remainder_clip_duration))

    # Concatenate the clips
    final_clip = concatenate_videoclips(repeated_clips)

    # Write the final video file with the proper FPS
    final_clip.write_videofile(output_video, fps=fps)
    final_clip.close()

def create_video_from_image(image_path, duration, output_path, fps=24):
    """
    Create a video from a single image that lasts for a specified duration.

    Parameters:
        image_path (str): Path to the image file.
        duration (float): Duration of the video in seconds.
        output_path (str): Output file path for the video.
        fps (int, optional): Frames per second of the output video. Default is 24.
    """
    # Read the image
    img = imageio.imread(image_path)

    # Create a writer object
    writer = imageio.get_writer(output_path, fps=fps)

    # Write the same image for the specified duration
    for _ in range(int(duration * fps)):
        writer.append_data(img)

    # Close the writer object
    writer.close()

def combine_videos(video_paths, output_path):
    """
    Combine multiple videos into one.

    Parameters:
        video_paths (list): List of paths to the videos to be combined.
        output_path (str): Output file path for the combined video.
    """
    # Load all video clips
    video_clips = [VideoFileClip(path) for path in video_paths]

    # Concatenate the video clips
    final_clip = concatenate_videoclips(video_clips)

    # Write the final video to the output path
    final_clip.write_videofile(output_path, codec="libx264")


def crop_to_circle(frame, circle_radius):
    """
    Crop the frame to a circle with the specified radius.

    Parameters:
        frame (numpy.ndarray): Input frame.
        circle_radius (int): Radius of the circle for cropping.

    Returns:
        numpy.ndarray: Cropped frame.
    """
    # Create a mask for the circle
    mask = np.zeros_like(frame)
    cv2.circle(mask, (frame.shape[1] // 2, frame.shape[0] // 2), circle_radius, (255, 255, 255), -1, 8, 0)

    # Apply the mask to the frame
    cropped_frame = cv2.bitwise_and(frame, mask)

    return cropped_frame

def overlay_video(background_path, overlay_path, circle_radius, output_path):
    # Load background video
    background_clip = VideoFileClip(background_path)

    # Load overlay video
    overlay_clip = VideoFileClip(overlay_path)

    # Crop overlay video to circle
    overlay_clip = overlay_clip.fl_image(lambda frame: crop_to_circle(frame, circle_radius))

    # Resize overlay video to fit in the bottom-right corner of the background video
    overlay_clip = overlay_clip.resize(width=int(background_clip.size[0] * 0.3))  # Adjust the scaling factor as needed

    # Set position of overlay video on background (bottom-right corner)
    overlay_clip = overlay_clip.set_position(("right", "bottom"))

    # Combine background and overlay videos
    final_clip = CompositeVideoClip([background_clip, overlay_clip])

    # Write the final video to the output path
    final_clip.write_videofile(output_path, codec="libx264", fps=background_clip.fps)
# API to create deepfake video using existing functionalities
@app.post("/deepfake/")
def create_deepfake(request: DeepFakeRequest):
    ppt_path = request.ppt_path
    image_path = request.image_path
    video_path = request.video_path

    # Extract speaker notes from PPT
    extracted_notes = extract_speaker_notes(ppt_path)

    # Convert notes to audio
    audio_paths = []
    for i, note in enumerate(extracted_notes):
        audio_path = comment_to_audio(note['speaker_notes'], i + 1)
        if audio_path:
            audio_paths.append(audio_path)

    source_image = imageio.imread(image_path)
    # Resize source image
    source_image = resize(source_image, (256, 256))[..., :3]

    # Initialize an empty list to store resized frames
    driving_video = []

    # Iterate over each frame in the video and resize it
    for frame in reader:
        frame = resize(frame, (256, 256))[..., :3]
        driving_video.append(frame)

    # Close the reader
    reader.close()

    predictions = demo.make_animation(source_image, driving_video, generator, kp_detector, relative=True)
    # Now you have source_image and driving_video available for further processing
    # Assuming predictions is defined and contains the generated frames
    save_generated_video(predictions, filename='./genratedvideo/predicted_video.mp4')

    duration = 0
    i = 1
    durations = []
    for note in extracted_notes:
        duration += audio_length(f'./outputaudio/output{i}.mp3')
        durations.append(audio_length(f'./outputaudio/output{i}.mp3'))
        i += 1
    time = duration / audio_length('./genratedvideo/predicted_video.mp4')

    input_audio_files = [ f"./outputaudio/output{x}.mp3" for x in range(i-1) ]

    output_audio_file = "./outputaudio/combined_audio.mp3"

    combine_audios(input_audio_files, output_audio_file)

    # Path to the input video
    input_video_path = "./genratedvideo/predicted_video.mp4"
    # Path to the output video
    output_video_path = "./genratedvideo/final_rep_video.mp4"

    repetitions = time

    repeat_video(input_video_path, repetitions, output_video_path)

    for j in range(i-1):
        image_path = f"./images/{j}.jpg"  # Replace with your image file path
        output_video_path = f"./genratedvideo/output_video{j}.mp4"
        duration = durations[j]
        create_video_from_image(image_path, duration, output_video_path, fps=24)
    # Dummy response for the sake of example

    video_paths = [f"./genratedvideo/output_video{x}.mp4"  for x  in range(i-1)]
    output_path = "./genratedvideo/combined_video.mp4"

    combine_videos(video_paths, output_path)

    # Paths to the input video and audio files
    video_path = "./genratedvideo/combined_video.mp4"
    audio_path = "./outputaudio/combined_audio.mp3"
    output_path = "./genratedvideo/combined_output_video.mp4"

    # Load the video and audio
    video_clip = VideoFileClip(video_path)
    audio_clip = AudioFileClip(audio_path)

    # Check if the video is shorter than the audio
    if video_clip.duration < audio_clip.duration:
        # Resize the video duration to match the audio
        video_clip = video_clip.set_duration(audio_clip.duration)

    # Set the audio of the video clip to the new audio
    final_clip = video_clip.set_audio(audio_clip)

    # Write the result to a file
    final_clip.write_videofile(output_path, codec='libx264', audio_codec='aac')

    # Close the clips
    video_clip.close()
    audio_clip.close()

    # Paths to the input video, audio, and output file
    video_path = "./genratedvideo/combined_video.mp4"
    audio_path = "./outputaudio/combined_audio.mp3"
    output_path = "./genratedvideo/combined_output_video.mp4"

    # Merge video and audio using ffmpeg
    ffmpeg_merge_video_audio(video=video_path, audio=audio_path, output=output_path)

    # Example usage:
    background_path = "./genratedvideo/combined_output_video.mp4"
    overlay_path = "./genratedvideo/final_rep_video.mp4"
    circle_radius = 500
    output_path = "./genratedvideo/final_video.mp4"
    overlay_video(background_path, overlay_path, circle_radius, output_path)

    shutil.rmtree('./genratedvideo')
    shutil.rmtree('./outputaudio')
    shutil.rmtree('./images')
    return {"message": "Deepfake creation started", "source_image": image_path, "video_path": video_path}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="127.0.0.1", port=8001)
